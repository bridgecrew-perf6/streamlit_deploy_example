from pptx import Presentation
from pptx.util import Inches, Pt,Mm
# you need to explictly make a letter size template to have align center working
#prs = Presentation("vcp_report_template.pptx")
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def fill_title_page_date():
    from streamlit_project_settings import VCP_REPORT_COVERPAGE_PATH
    #prs = Presentation("vcp_report_template.pptx")
    prs = Presentation(VCP_REPORT_COVERPAGE_PATH)
    #slide_layout = prs.slide_layouts[0]  # assuming you want the first one
    slide = prs.slides[0]
    from datetime import datetime
    for shape in slide.shapes:
        if shape.name == 'date_slot': # this refers to the date
            text_frame = shape.text_frame
            #import tzlocal
            import pytz
            from datetime import datetime as dt_
            now = str(dt_.now(pytz.timezone("US/Pacific")).strftime('%Y-%m-%d %H:%M:%S %z'))
            #now = datetime.now(tzlocal.get_localzone())
            #text_frame.text = str(now)
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text =  str(now)
            font = run.font
            from pptx.dml.color import RGBColor
            font.color.rgb = RGBColor(255, 255, 255)

    return prs

def fill_table(df, table):
    import numpy as np
    dfvals = df.values
    # d
    for n_col in np.arange(0, len(df.columns)):
        cell = table.cell(0, n_col)
        cell.text = str(df.columns[n_col])

    for idx in np.arange(0, df.shape[0]):
        idx_aug = idx + 1
        for idy in np.arange(0, df.shape[1]):
            cell = table.cell(idx_aug, idy)
            #print(idx, idy)
            cell.text = str(dfvals[idx, idy])

    # these functions are change font
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
    iter_cells(table)
    pass


def build_powerpoint_page(ticker, page_number, time_series_png_path,vcp_property_df,stockinfo_df,factors_df, vcp_supplementarty_df,contractions_df):
    from streamlit_project_settings import SINGLE_POWERPOINT_PAGE_PATH
    prs = Presentation(SINGLE_POWERPOINT_PAGE_PATH)
    #prs = Presentation(
    #    r"C:\Users\tclyu\PycharmProjects\streamlit_deploy_example\pdf_report_generator\page_template.pptx")
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

    # put down title and page number
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.name == 'page_title': # this refers to the date
            text_frame = shape.text_frame
            text_frame.text = str(ticker)
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
        if shape.name == 'page_number':
            text_frame = shape.text_frame
            text_frame.text = str(page_number)
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)

    placeholder = slide.shapes#placeholders[1]
    placeholder.alignment = PP_PARAGRAPH_ALIGNMENT

    TIME_SERIES_TOP = 0.4 # inches
    TIME_SERIES_IMAGE_HEIGHT = 3. #inches

    ######################## put time series plot ##########################
    #https://stackoverflow.com/questions/40149992/python-pptx-align-image-to-center-of-slide
    picture = placeholder.add_picture(time_series_png_path,left=False,top=Inches(TIME_SERIES_TOP),
                                      width=Inches(7), height=Inches(TIME_SERIES_IMAGE_HEIGHT))
    picture.left = int((prs.slide_width - picture.width) / 2)
    #########################################################################
    VCP_PROPERTY_TOP = TIME_SERIES_IMAGE_HEIGHT + TIME_SERIES_TOP
    # ---add table to slide---
    # x, y, cx, cy = Inches(0.25), Inches(VCP_PROPERTY_TOP), Inches(7), Inches(1)
    x, y, cx, cy = Inches(0.25), Inches(VCP_PROPERTY_TOP), Inches(3.5), Inches(1)
    shape = slide.shapes.add_table(vcp_property_df.shape[0] + 1, vcp_property_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(vcp_property_df, table)
    next_top = y  # +2.2*cy
    #print(next_top)
    # ---add table to slide---
    # x, y, cx, cy = Inches(0.25), (next_top), Inches(7), Inches(1.)
    x, y, cx, cy = x + cx, (next_top), Inches(3.5), Inches(1.)
    shape = slide.shapes.add_table(stockinfo_df.shape[0] + 1, stockinfo_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(stockinfo_df, table)
    next_top = next_top + 2.3 * cy
    # ---add table to slide---
    x, y, cx, cy = Inches(0.25), next_top, Inches(7), Inches(1.)
    shape = slide.shapes.add_table(factors_df.shape[0] + 1, factors_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(factors_df, table)

    next_top = next_top + 1.15 * cy
    # ---add table to slide---
    x, y, cx, cy = Inches(0.25), next_top, Inches(7), Inches(1.)
    shape = slide.shapes.add_table(vcp_supplementarty_df.shape[0] + 1, vcp_supplementarty_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(vcp_supplementarty_df, table)

    next_top = next_top + 1.75 * cy
    # ---add table to slide---
    x, y, cx, cy = Inches(0.25), next_top, Inches(7), Inches(1.)
    shape = slide.shapes.add_table(contractions_df.shape[0] + 1, contractions_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(contractions_df, table)
    return prs

def build_content_page(selected_df):
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches
    #selected_df = pd.read_pickle(os.path.join(dir, "SELECTEDSUMMARY_selected_df.pkl"))

    from streamlit_project_settings import SINGLE_POWERPOINT_PAGE_PATH
    prs = Presentation(SINGLE_POWERPOINT_PAGE_PATH)
    #prs = Presentation(
    #    r"C:\Users\tclyu\PycharmProjects\streamlit_deploy_example\pdf_report_generator\page_template.pptx")
    slide = prs.slides[0]  # .add_slide(prs.slide_layouts[5])

    x, y, cx, cy = Inches(0.25), Inches(1), Inches(7), Inches(5.)
    shape = slide.shapes.add_table(selected_df.shape[0] + 1, selected_df.shape[1], x, y, cx, cy)
    table = shape.table
    fill_table(selected_df, table)
    return prs

def build_pdf_report_from_dir(dir, output_pdf_path, ppt_export_dir = "tmp_ppt_export", is_cleanup=True):
    import os
    """
    def get_tickers(dir):
        import os
        tickers = []
        for file in os.listdir(dir):
            if "SELECTEDSUMMARY" not in file:
                if ".png" in file:
                    tmp_filename = file.replace('.png', '')
                if ".html" in file:
                    tmp_filename = file.replace('.html', '')
                tickers.append(tmp_filename.split('_')[0])
        tickers = sorted(list(set(tickers)))
        return tickers
    ticker_list = get_tickers(dir)
    """

    ppt_export_dir = os.path.abspath(ppt_export_dir)
    if os.path.exists(ppt_export_dir):
        import shutil
        if os.path.exists(ppt_export_dir):
            shutil.rmtree(ppt_export_dir, ignore_errors=False, onerror=None)

    from pathlib import Path
    import pandas as pd

    Path(ppt_export_dir).mkdir(exist_ok=True)
    from pdf_report_generator.powerpoint_report_gen_utils.utils import PPTtoPDF, merge_pdfs

    pdfs_to_combine = []

    ######### title page #######
    prs = fill_title_page_date()
    inputFileName = os.path.join(ppt_export_dir, 'tmp.pptx')
    outputFileName = os.path.join(ppt_export_dir, "title.pdf")
    prs.save(inputFileName)
    PPTtoPDF(inputFileName, outputFileName, formatType=32)
    pdfs_to_combine.append(os.path.abspath((outputFileName)))

    ######### content page #######
    selected_df = pd.read_pickle(os.path.join(dir, "SELECTEDSUMMARY_selected_df.pkl"))
    prs = build_content_page(selected_df)
    inputFileName = os.path.join(ppt_export_dir, 'tmp.pptx')
    outputFileName = os.path.join(ppt_export_dir, "contents.pdf")
    prs.save(inputFileName)
    PPTtoPDF(inputFileName, outputFileName, formatType=32)
    pdfs_to_combine.append(os.path.abspath((outputFileName)))

    for idx, row in selected_df.iterrows():
        ticker = row.ticker
        page_number = row.page

        time_series_png_path = os.path.join(dir, ticker + '.png')
        vcp_property_df = pd.read_pickle(os.path.join(dir, ticker + "_vcp_property_df.pkl"))
        stockinfo_df = pd.read_pickle(os.path.join(dir, ticker + "_stockinfo_df.pkl"))
        factors_df = pd.read_pickle(os.path.join(dir, ticker + "_factors_df.pkl"))
        vcp_supplementarty_df = pd.read_pickle(os.path.join(dir, ticker + "_vcp_supplementary_df.pkl"))
        contractions_df = pd.read_pickle(os.path.join(dir, ticker + "_contractions_df.pkl"))

        prs = build_powerpoint_page(ticker, page_number, time_series_png_path, vcp_property_df, stockinfo_df,
                                    factors_df,
                                    vcp_supplementarty_df, contractions_df)
        inputFileName = os.path.join(ppt_export_dir, 'tmp.pptx')
        outputFileName = os.path.join(ppt_export_dir, str(page_number) + '_' + str(ticker) + ".pdf")
        prs.save(inputFileName)

        PPTtoPDF(inputFileName, outputFileName, formatType=32)
        print("build_pdf_report_from_dir::converted ppt to pdf=",outputFileName)
        pdfs_to_combine.append(os.path.abspath((outputFileName)))

    merge_pdfs(pdfs_to_combine, output_pdf_path)

    if is_cleanup and os.path.exists(ppt_export_dir):
        import shutil
        if os.path.exists(ppt_export_dir):
            shutil.rmtree(ppt_export_dir, ignore_errors=False, onerror=None)

    print("build_pdf_report_from_dir:: output= "+os.path.abspath(output_pdf_path))

    return output_pdf_path

if __name__ == "____":
    import pandas as pd
    dir=r"C:\Users\tclyu\PycharmProjects\streamlit_deploy_example\report_data_export_2022-04-28"
    import os
    df = pd.read_pickle(os.path.join(dir,"SELECTEDSUMMARY_selected_df.pkl"))
    print(df)
if __name__ == "__main__":

    prs = fill_title_page_date()
    prs.save('tmp.pptx')